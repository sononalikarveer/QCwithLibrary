import re
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import RGBColor

def clean_text(text):
    return re.sub(r"\s+", " ", text).strip()

# =========================
# SLIDE TEXT EXTRACTION
# =========================
def get_slide_text(slide):
    points = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if not text:
                continue

            font_rgb = None
            if (
                para.runs
                and para.runs[0].font.color
                and para.runs[0].font.color.type == MSO_COLOR_TYPE.RGB
            ):
                font_rgb = para.runs[0].font.color.rgb

            # Ignore decorative orange text
            if font_rgb != RGBColor(242, 103, 34):
                points.append(clean_text(text))

    return points

# =========================
# VO EXTRACTION
# =========================
def get_vo_text(slide):
    if not slide.has_notes_slide:
        return []

    notes = slide.notes_slide.notes_text_frame.text

    match = re.search(
        r"(?i)vo:\s*(.*?)(Image Link:|Instructions to GD:|$)",
        notes,
        re.DOTALL
    )

    if not match:
        return []

    vo_block = match.group(1).strip()

    return [
        clean_text(line.strip("-• "))
        for line in vo_block.split("\n")
        if line.strip()
    ]
